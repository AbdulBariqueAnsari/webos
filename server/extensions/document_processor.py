import os, json, re, html, base64, mimetypes
from datetime import datetime

class DocumentProcessor:
    def __init__(self, storage_dir):
        self.storage_dir = storage_dir
        self.supported = {
            ".txt": "text/plain",
            ".md": "text/markdown",
            ".html": "text/html",
            ".json": "application/json",
            ".xml": "application/xml",
            ".csv": "text/csv",
            ".py": "text/x-python",
            ".js": "text/javascript",
            ".css": "text/css",
            ".ini": "text/plain",
            ".cfg": "text/plain",
            ".log": "text/plain",
            ".yaml": "text/yaml",
            ".yml": "text/yaml",
            ".pdf": "application/pdf",
        }

    def read(self, path):
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            return self._read_pdf(path)
        elif ext in (".docx", ".doc"):
            return self._read_docx(path)
        elif ext in (".xlsx", ".xls"):
            return self._read_xlsx(path)
        elif ext == ".pptx":
            return self._read_pptx(path)
        elif ext in self.supported:
            return self._read_text(path)
        else:
            return self._read_binary(path)

    def ls(self, dir_path):
        results = []
        try:
            for f in os.listdir(dir_path):
                fp = os.path.join(dir_path, f)
                stat = os.stat(fp)
                ext = os.path.splitext(f)[1].lower()
                results.append({
                    "name": f,
                    "path": fp,
                    "size": stat.st_size,
                    "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    "type": "dir" if os.path.isdir(fp) else "file",
                    "mime": self.supported.get(ext, "application/octet-stream"),
                    "ext": ext
                })
        except Exception as e:
            return {"error": str(e)}
        return sorted(results, key=lambda x: (-1 if x["type"]=="dir" else 1, x["name"]))

    def search(self, root, query):
        results = []
        query_lower = query.lower()
        for root_dir, dirs, files in os.walk(root):
            for f in files:
                if query_lower in f.lower():
                    fp = os.path.join(root_dir, f)
                    try:
                        results.append({"name": f, "path": fp, "size": os.path.getsize(fp)})
                    except Exception:
                        pass
                if len(results) >= 50:
                    break
            if len(results) >= 50:
                break
        return results

    def _read_text(self, path):
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            ext = os.path.splitext(path)[1].lower()
            return {"type": "text", "content": content, "mime": self.supported.get(ext, "text/plain"), "size": len(content)}
        except Exception as e:
            return {"error": str(e)}

    def _read_binary(self, path):
        try:
            with open(path, "rb") as f:
                data = f.read()
            ext = os.path.splitext(path)[1].lower()
            mime = self.supported.get(ext, "application/octet-stream")
            b64 = base64.b64encode(data).decode()
            return {"type": "binary", "mime": mime, "size": len(data), "base64": b64[:100000]}
        except Exception as e:
            return {"error": str(e)}

    def _read_pdf(self, path):
        try:
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                pages = []
                for i, page in enumerate(reader.pages[:50]):
                    pages.append({"page": i+1, "text": page.extract_text() or ""})
                return {"type": "pdf", "pages": len(pages), "content": pages, "meta": {"pages": len(reader.pages)}}
        except ImportError:
            return self._read_binary(path)
        except Exception as e:
            return {"type": "pdf", "error": str(e), "fallback": self._read_binary(path)}

    def _read_docx(self, path):
        try:
            try:
                from docx import Document
                doc = Document(path)
                paragraphs = [p.text for p in doc.paragraphs]
                tables = []
                for table in doc.tables:
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text for cell in row.cells])
                    tables.append(rows)
                return {"type": "docx", "paragraphs": paragraphs[:200], "tables": tables[:10], "total_paragraphs": len(paragraphs)}
            except ImportError:
                import zipfile
                import xml.etree.ElementTree as ET
                with zipfile.ZipFile(path) as z:
                    xml = z.read("word/document.xml")
                    root = ET.fromstring(xml)
                    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                    texts = []
                    for t in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                        if t.text:
                            texts.append(t.text)
                    return {"type": "docx", "paragraphs": texts[:200], "note": "Basic XML extraction"}
        except Exception as e:
            return {"error": str(e)}

    def _read_xlsx(self, path):
        try:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                sheets = []
                for name in wb.sheetnames[:10]:
                    ws = wb[name]
                    rows = []
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i > 500:
                            break
                        rows.append([str(c) if c is not None else "" for c in row])
                    sheets.append({"name": name, "rows": rows, "total_rows": ws.max_row})
                wb.close()
                return {"type": "xlsx", "sheets": sheets}
            except ImportError:
                import zipfile, xml.etree.ElementTree as ET
                with zipfile.ZipFile(path) as z:
                    sheets = []
                    i = 0
                    while True:
                        name = f"xl/worksheets/sheet{i+1}.xml"
                        if name not in z.namelist():
                            break
                        xml = z.read(name)
                        root = ET.fromstring(xml)
                        ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                        rows_data = []
                        for row in root.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}row"):
                            cells = []
                            for c in row.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}c"):
                                v = c.find("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}v")
                                cells.append(v.text if v is not None else "")
                            rows_data.append(cells)
                        sheets.append({"name": f"Sheet{i+1}", "rows": rows_data[:200]})
                        i += 1
                    return {"type": "xlsx", "sheets": sheets}
        except Exception as e:
            return {"error": str(e)}

    def _read_pptx(self, path):
        try:
            try:
                from pptx import Presentation
                prs = Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides[:30]):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                    slides.append({"slide": i+1, "texts": texts})
                return {"type": "pptx", "slides": slides}
            except ImportError:
                return self._read_binary(path)
        except Exception as e:
            return {"error": str(e)}

    def preview(self, path, max_size=100000):
        try:
            size = os.path.getsize(path)
            if size > max_size * 10:
                return {"type": "too_large", "size": size, "max_size": max_size}
            return self.read(path)
        except Exception as e:
            return {"error": str(e)}
